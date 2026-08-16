"""
Extract text from the official M110 slide decks.
استخراج النص من شرائح مقرر M110 الرسمية

Walks slides-official/, reads every .pptx, and writes a sibling .pptx.txt
holding the deck's text in slide order.
يمر على مجلد الشرائح ويكتب ملف نصي بجانب كل ملف عرض تقديمي.

These text files let AI assistants and text search read the official slides.
They do NOT replace the PDFs: figures, flowcharts and diagrams are images and
do not survive extraction.
هذه الملفات النصية لا تُغني عن ملفات PDF: الأشكال والمخططات صور ولا تُستخرج.

Usage / الاستخدام:
    pip install python-pptx
    python tools/extract_slides.py           # write the .txt files
    python tools/extract_slides.py --check   # compare only, write nothing

Note on --check results / ملاحظة حول نتائج --check:
This script does not reproduce the committed .txt files byte-for-byte.
Body text matches closely (mean word-similarity ~0.958, measured by
comparing whitespace-split tokens against the committed files). Two
verified, minor structural differences remain: the committed files
contain more blank lines (the "if text.strip():" guard below skips
empty paragraphs that the original conversion kept), and a handful of
slide numbers appear more than once in the committed files but only once
here -- the original conversion evidently walked the slide shapes
differently; the exact mechanism is not known. Field text such as slide
numbers and dates is recovered normally by this script. Regenerate only
if the decks themselves are replaced, and expect these small structural
differences.
لا يُعيد هذا السكربت إنتاج الملفات النصية المرفوعة حرفاً بحرف. يتطابق
النص الأساسي بدرجة عالية (متوسط تشابه الكلمات ≈0.958، مقارنة بالملفات
المرفوعة). يبقى فرقان بنيويان بسيطان تم التحقق منهما: تحتوي الملفات
المرفوعة على أسطر فارغة أكثر (الشرط "if text.strip():" أدناه يتجاهل
الفقرات الفارغة التي احتفظ بها التحويل الأصلي)، وتظهر بعض أرقام الشرائح
أكثر من مرة في الملفات المرفوعة بينما تظهر مرة واحدة فقط هنا -- يبدو أن
التحويل الأصلي تعامل مع عناصر الشريحة بطريقة مختلفة، لكن الآلية الدقيقة
غير معروفة. يتم استرجاع نص الحقول مثل أرقام الشرائح والتواريخ بشكل طبيعي
بواسطة هذا السكربت. أعد التوليد فقط عند استبدال الشرائح نفسها، وتوقّع
وجود هذه الفروق البنيوية الصغيرة.
"""

import argparse
import sys
from pathlib import Path

from pptx import Presentation

# slides-official/ sits next to tools/ / مجلد الشرائح بجانب مجلد الأدوات
SLIDES_DIR = Path(__file__).resolve().parent.parent / "slides-official"


def extract_deck(pptx_path):
    """
    Return every slide's text, blank line between slides.
    يُرجع نص كل شريحة، مع سطر فارغ بين الشرائح.
    """
    presentation = Presentation(str(pptx_path))
    blocks = []

    for slide in presentation.slides:
        lines = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text
                if text.strip():
                    lines.append(text)
        blocks.append("\n".join(lines))

    return "\n\n".join(blocks) + "\n"


def main():
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--check",
        action="store_true",
        help="Compare against existing .txt files without writing",
    )
    args = parser.parse_args()

    decks = sorted(SLIDES_DIR.glob("*/*.pptx"))
    if not decks:
        print(f"No .pptx files found under {SLIDES_DIR}")
        return 1

    differences = 0
    for deck in decks:
        target = deck.with_suffix(".pptx.txt")
        extracted = extract_deck(deck)

        if args.check:
            existing = target.read_text(encoding="utf-8") if target.exists() else ""
            status = "same" if existing == extracted else "DIFFERS"
            differences += status == "DIFFERS"
            print(f"{status:8} {target.relative_to(SLIDES_DIR.parent)}")
        else:
            target.write_text(extracted, encoding="utf-8")
            print(f"wrote    {target.relative_to(SLIDES_DIR.parent)}")

    return 1 if differences else 0


if __name__ == "__main__":
    sys.exit(main())
